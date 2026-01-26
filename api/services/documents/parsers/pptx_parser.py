# api/services/documents/parsers/pptx_parser.py
"""
Parser for Microsoft PowerPoint presentations (.pptx, .ppt).
"""

import logging
import subprocess
import tempfile
from pathlib import Path
from typing import BinaryIO, Optional, Union
import time

from ..base import (
    BaseParser, DocumentContent, DocumentMetadata,
    FileType, ParserResult, TextChunk
)

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Parser for PowerPoint presentations."""
    
    supported_types = [FileType.PPTX, FileType.PPT]
    
    async def parse(
        self,
        file_path: Union[str, Path],
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> ParserResult:
        """Parse a PowerPoint presentation from file path."""
        start_time = time.time()
        file_path = Path(file_path)
        
        try:
            with open(file_path, "rb") as f:
                file_data = f.read()
            
            result = await self._parse_content(
                file_data,
                file_path.name,
                chunk_size,
                chunk_overlap,
            )
            result.processing_time_ms = (time.time() - start_time) * 1000
            return result
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint {file_path}: {e}")
            return ParserResult(
                success=False,
                error=str(e),
                processing_time_ms=(time.time() - start_time) * 1000,
            )
    
    async def parse_bytes(
        self,
        file_data: BinaryIO,
        filename: str,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> ParserResult:
        """Parse a PowerPoint presentation from bytes."""
        start_time = time.time()
        
        try:
            data = file_data.read()
            result = await self._parse_content(
                data,
                filename,
                chunk_size,
                chunk_overlap,
            )
            result.processing_time_ms = (time.time() - start_time) * 1000
            return result
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint bytes: {e}")
            return ParserResult(
                success=False,
                error=str(e),
                processing_time_ms=(time.time() - start_time) * 1000,
            )
    
    async def _parse_content(
        self,
        file_data: bytes,
        filename: str,
        chunk_size: Optional[int],
        chunk_overlap: Optional[int],
    ) -> ParserResult:
        """Parse PowerPoint content."""
        file_type = self.detect_file_type(filename)
        file_hash = self.compute_file_hash(file_data)
        warnings = []
        
        with tempfile.NamedTemporaryFile(
            suffix=f".{file_type.value}", delete=False
        ) as tmp:
            tmp.write(file_data)
            tmp_path = Path(tmp.name)
        
        try:
            # Convert .ppt to .pptx if needed
            if file_type == FileType.PPT:
                converted_path = await self._convert_ppt_to_pptx(tmp_path)
                if converted_path:
                    tmp_path = converted_path
                else:
                    warnings.append("Could not convert .ppt to .pptx")
                    return ParserResult(
                        success=False,
                        error="Cannot parse .ppt format without LibreOffice",
                        warnings=warnings,
                    )
            
            # Extract content
            slides_text, chunks, tables = await self._extract_with_pptx(
                tmp_path, chunk_size, chunk_overlap
            )
            
            # Extract metadata
            metadata = await self._extract_metadata(
                tmp_path, file_data, filename, file_type, file_hash
            )
            
            content = DocumentContent(
                full_text=slides_text,
                chunks=chunks,
                tables=tables,
            )
            
            return ParserResult(
                success=True,
                metadata=metadata,
                content=content,
                warnings=warnings,
            )
            
        finally:
            tmp_path.unlink(missing_ok=True)
            if file_type == FileType.PPT:
                pptx_path = tmp_path.with_suffix(".pptx")
                pptx_path.unlink(missing_ok=True)
    
    async def _convert_ppt_to_pptx(self, ppt_path: Path) -> Optional[Path]:
        """Convert .ppt to .pptx using LibreOffice."""
        try:
            output_dir = ppt_path.parent
            result = subprocess.run(
                [
                    "soffice", "--headless", "--convert-to", "pptx",
                    "--outdir", str(output_dir), str(ppt_path)
                ],
                capture_output=True,
                timeout=120,
            )
            if result.returncode == 0:
                pptx_path = ppt_path.with_suffix(".pptx")
                if pptx_path.exists():
                    return pptx_path
        except Exception as e:
            logger.warning(f"Failed to convert .ppt to .pptx: {e}")
        return None
    
    async def _extract_with_pptx(
        self,
        file_path: Path,
        chunk_size: Optional[int],
        chunk_overlap: Optional[int],
    ) -> tuple[str, list[TextChunk], list[dict]]:
        """Extract content using python-pptx library."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.error("python-pptx not installed")
            return "", [], []
        
        try:
            prs = Presentation(str(file_path))
            
            all_text = []
            chunks = []
            tables = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                
                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text_parts.append(text)
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = self._extract_table(shape.table)
                        if table_data:
                            tables.append({
                                "slide": slide_num,
                                "rows": table_data,
                            })
                            # Also add table content as text
                            for row in table_data:
                                slide_text_parts.append(" | ".join(row))
                
                if slide_text_parts:
                    slide_content = "\n".join(slide_text_parts)
                    all_text.append(f"[Slide {slide_num}]\n{slide_content}")
                    
                    # Create chunk for this slide
                    chunks.append(TextChunk(
                        content=slide_content,
                        source=f"slide_{slide_num}",
                        chunk_index=slide_num - 1,
                        page_number=slide_num,
                    ))
            
            full_text = "\n\n".join(all_text)
            
            # Re-chunk if needed
            if chunk_size and any(len(c.content) > chunk_size for c in chunks):
                chunks = self.chunk_text(full_text, "chunk", chunk_size, chunk_overlap)
            
            return full_text, chunks, tables
            
        except Exception as e:
            logger.error(f"Error with python-pptx: {e}")
            return "", [], []
    
    def _extract_table(self, table) -> list[list[str]]:
        """Extract table data."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text.strip() if cell.text else ""
                cells.append(text)
            rows.append(cells)
        return rows
    
    async def _extract_metadata(
        self,
        file_path: Path,
        file_data: bytes,
        filename: str,
        file_type: FileType,
        file_hash: str,
    ) -> DocumentMetadata:
        """Extract presentation metadata."""
        metadata = DocumentMetadata(
            filename=filename,
            file_type=file_type,
            file_size=len(file_data),
            file_hash=file_hash,
        )
        
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            
            metadata.slide_count = len(prs.slides)
            
            # Core properties
            core_props = prs.core_properties
            metadata.title = core_props.title
            metadata.author = core_props.author
            metadata.subject = core_props.subject
            metadata.created_at = core_props.created
            metadata.modified_at = core_props.modified
            
        except Exception as e:
            logger.warning(f"Could not extract PowerPoint metadata: {e}")
        
        return metadata
